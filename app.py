import streamlit as st
import tempfile
import os
import random
from topic_extractor import extract_topics
from quiz_generator import generate_quiz
from groq import Groq
from dotenv import load_dotenv
import fitz
from pptx import Presentation

load_dotenv()
client = Groq(api_key=os.getenv("GROQ_API_KEY"))

st.set_page_config(
    page_title="AI Learning Coach",
    page_icon="🧠",
    layout="wide"
)

QUOTES = [
    "🔥 Success is the sum of small efforts repeated daily!",
    "💡 The expert in anything was once a beginner!",
    "🚀 Push yourself, no one else is going to do it for you!",
    "⚡ Dream it. Wish it. Do it!",
    "🎯 Focus on your goal, distractions are temporary!",
    "💪 Hard work beats talent when talent doesn't work hard!",
    "🌟 Believe you can and you're halfway there!",
    "📚 Education is the passport to the future!",
    "🏆 Winners are not people who never fail, but never quit!",
    "✨ Your only limit is your mind!",
]

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700;800&display=swap');
* { font-family: 'Inter', sans-serif; }
.stApp { background: #0a0a0f; }
.hero { text-align: center; padding: 3rem 0 2rem 0; }
.hero-tag {
    display: inline-block;
    background: linear-gradient(90deg, #a78bfa22, #60a5fa22);
    border: 1px solid #a78bfa55;
    color: #a78bfa;
    padding: 0.3rem 1rem;
    border-radius: 20px;
    font-size: 0.8rem;
    letter-spacing: 2px;
    text-transform: uppercase;
    margin-bottom: 1rem;
}
.hero-title {
    font-size: 4rem;
    font-weight: 800;
    background: linear-gradient(135deg, #ffffff 0%, #a78bfa 50%, #60a5fa 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    line-height: 1.1;
    margin-bottom: 1rem;
}
.hero-sub {
    color: #64748b;
    font-size: 1.1rem;
    max-width: 500px;
    margin: 0 auto 2rem auto;
}
.quote-box {
    background: linear-gradient(135deg, #a78bfa11, #60a5fa11);
    border: 1px solid #a78bfa33;
    border-radius: 16px;
    padding: 1.2rem 1.8rem;
    margin-bottom: 1rem;
    text-align: center;
    color: #c4b5fd;
    font-size: 1rem;
    font-weight: 500;
    font-style: italic;
}
.stat-row {
    display: flex;
    gap: 1rem;
    margin: 2rem 0;
}
.stat-card {
    flex: 1;
    background: rgba(255,255,255,0.02);
    border: 1px solid rgba(255,255,255,0.06);
    border-radius: 16px;
    padding: 1.2rem;
    text-align: center;
}
.stat-num {
    font-size: 2rem;
    font-weight: 800;
    background: linear-gradient(90deg, #a78bfa, #60a5fa);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
}
.stat-lbl { color: #475569; font-size: 0.8rem; margin-top: 0.2rem; }
.glass-card {
    background: rgba(255,255,255,0.03);
    border: 1px solid rgba(255,255,255,0.08);
    border-radius: 20px;
    padding: 1.8rem;
    margin: 1rem 0;
}
.tips-card {
    background: rgba(255,255,255,0.02);
    border: 1px solid rgba(255,255,255,0.06);
    border-radius: 16px;
    padding: 1.5rem;
    margin-top: 1rem;
}
.tip-item {
    display: flex;
    align-items: flex-start;
    gap: 0.8rem;
    padding: 0.6rem 0;
    border-bottom: 1px solid rgba(255,255,255,0.04);
    color: #94a3b8;
    font-size: 0.85rem;
    line-height: 1.5;
}
.tip-item:last-child { border-bottom: none; }
.badge {
    display: inline-block;
    background: linear-gradient(135deg, #a78bfa15, #60a5fa15);
    border: 1px solid #a78bfa33;
    color: #c4b5fd;
    padding: 0.4rem 1rem;
    border-radius: 50px;
    margin: 0.3rem;
    font-size: 0.82rem;
    font-weight: 500;
}
.q-card {
    background: rgba(255,255,255,0.02);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 16px;
    padding: 1.5rem;
    margin: 1rem 0;
}
.q-number {
    color: #a78bfa;
    font-size: 0.8rem;
    font-weight: 600;
    letter-spacing: 1px;
    text-transform: uppercase;
    margin-bottom: 0.5rem;
}
.q-text {
    color: #e2e8f0;
    font-size: 1rem;
    font-weight: 600;
    margin-bottom: 1rem;
    line-height: 1.5;
}
.correct {
    background: rgba(52,211,153,0.1) !important;
    border: 1px solid #34d39944 !important;
    border-radius: 12px !important;
    color: #34d399 !important;
}
.wrong {
    background: rgba(248,113,113,0.1) !important;
    border: 1px solid #f8717144 !important;
    border-radius: 12px !important;
    color: #f87171 !important;
}
.score-box {
    background: linear-gradient(135deg, #a78bfa22, #60a5fa22);
    border: 1px solid #a78bfa44;
    border-radius: 20px;
    padding: 2rem;
    text-align: center;
    margin: 2rem 0;
}
.score-num {
    font-size: 3rem;
    font-weight: 800;
    background: linear-gradient(90deg, #a78bfa, #60a5fa);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
}
.chat-user {
    background: rgba(167,139,250,0.1);
    border: 1px solid #a78bfa33;
    border-radius: 12px;
    padding: 1rem;
    margin: 0.5rem 0;
    color: #e2e8f0;
}
.chat-ai {
    background: rgba(96,165,250,0.1);
    border: 1px solid #60a5fa33;
    border-radius: 12px;
    padding: 1rem;
    margin: 0.5rem 0;
    color: #e2e8f0;
}
.stButton > button {
    background: linear-gradient(135deg, #a78bfa, #60a5fa) !important;
    color: white !important;
    border: none !important;
    border-radius: 14px !important;
    padding: 0.8rem 2rem !important;
    font-size: 1rem !important;
    font-weight: 700 !important;
    width: 100% !important;
    transition: all 0.3s !important;
}
.stButton > button:hover {
    transform: translateY(-3px) !important;
    box-shadow: 0 12px 30px rgba(167,139,250,0.35) !important;
}
.divider {
    height: 1px;
    background: linear-gradient(90deg, #a78bfa33, transparent);
    margin: 2rem 0;
}
.section-title {
    font-size: 1.4rem;
    font-weight: 700;
    color: #e2e8f0;
    margin: 2rem 0 1rem 0;
}
div[role="radiogroup"] label {
    background: rgba(255,255,255,0.03) !important;
    border: 1px solid rgba(255,255,255,0.08) !important;
    border-radius: 10px !important;
    padding: 0.6rem 1rem !important;
    margin: 0.3rem 0 !important;
    color: #94a3b8 !important;
    width: 100% !important;
}
div[role="radiogroup"] label:hover {
    border-color: #a78bfa44 !important;
    background: rgba(167,139,250,0.05) !important;
}
::-webkit-scrollbar { width: 6px; }
::-webkit-scrollbar-track { background: #0a0a0f; }
::-webkit-scrollbar-thumb { background: #a78bfa44; border-radius: 3px; }
</style>
""", unsafe_allow_html=True)


def read_uploaded_file(uploaded_file):
    text = ""
    if uploaded_file.name.endswith(".txt"):
        text = uploaded_file.read().decode("utf-8")
    elif uploaded_file.name.endswith(".pdf"):
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as f:
            f.write(uploaded_file.read())
            temp_path = f.name
        doc = fitz.open(temp_path)
        for page in doc:
            text += page.get_text()
        os.unlink(temp_path)
    elif uploaded_file.name.endswith(".pptx"):
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as f:
            f.write(uploaded_file.read())
            temp_path = f.name
        prs = Presentation(temp_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text_frame.text + "\n"
        os.unlink(temp_path)
    return text


def parse_quiz(quiz_text):
    questions = []
    current_q = {}
    for line in quiz_text.strip().split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("Q") and "." in line[:4]:
            if current_q:
                questions.append(current_q)
            current_q = {"question": line, "options": [], "answer": ""}
        elif line.startswith(("A.", "B.", "C.", "D.")):
            if current_q:
                current_q["options"].append(line)
        elif line.lower().startswith("answer"):
            if current_q:
                current_q["answer"] = line.split(":")[-1].strip()
    if current_q:
        questions.append(current_q)
    return questions


# Session State
if "questions_data" not in st.session_state:
    st.session_state.questions_data = []
if "user_answers" not in st.session_state:
    st.session_state.user_answers = {}
if "submitted" not in st.session_state:
    st.session_state.submitted = False
if "quote" not in st.session_state:
    st.session_state.quote = random.choice(QUOTES)
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# HERO
st.markdown("""
<div class="hero">
    <div class="hero-tag">✨ AI Powered</div>
    <div class="hero-title">AI Learning Coach</div>
    <div class="hero-sub">Upload your syllabus and get personalized quizzes instantly</div>
</div>
""", unsafe_allow_html=True)

# QUOTE
st.markdown(f'<div class="quote-box">"{st.session_state.quote}"</div>', unsafe_allow_html=True)

# STATS
st.markdown("""
<div class="stat-row">
    <div class="stat-card"><div class="stat-num">PDF</div><div class="stat-lbl">TXT • PPTX • PDF</div></div>
    <div class="stat-card"><div class="stat-num">50+</div><div class="stat-lbl">Questions Per Topic</div></div>
    <div class="stat-card"><div class="stat-num">AI</div><div class="stat-lbl">Groq Powered</div></div>
    <div class="stat-card"><div class="stat-num">⚡</div><div class="stat-lbl">Instant Results</div></div>
</div>
""", unsafe_allow_html=True)

st.markdown('<div class="divider"></div>', unsafe_allow_html=True)

# UPLOAD
left, right = st.columns([1.2, 1], gap="large")

with left:
    st.markdown('<div class="glass-card">', unsafe_allow_html=True)
    st.markdown("### 📂 Upload Karo")
    uploaded_file = st.file_uploader(
        "PDF, TXT ya PPTX",
        type=["txt", "pdf", "pptx"],
        label_visibility="collapsed"
    )
    if uploaded_file:
        st.success(f"✅ **{uploaded_file.name}** ready!")
    st.markdown("<br>", unsafe_allow_html=True)
    generate_btn = st.button("⚡ Generate Topics & Quiz")
    st.markdown('</div>', unsafe_allow_html=True)

with right:
    st.markdown("""
    <div class="glass-card">
        <div style="color:#e2e8f0; font-weight:700; font-size:1.1rem; margin-bottom:1rem;">💡 Kaise Kaam Karta Hai?</div>
        <div style="color:#64748b; font-size:0.9rem; line-height:2.2;">
            📂 &nbsp;<span style="color:#94a3b8">Apna syllabus upload karo</span><br>
            🔍 &nbsp;<span style="color:#94a3b8">AI topics identify karega</span><br>
            ❓ &nbsp;<span style="color:#94a3b8">Quiz attempt karo</span><br>
            📊 &nbsp;<span style="color:#94a3b8">Score check karo</span><br>
            💬 &nbsp;<span style="color:#94a3b8">Doubts poochho AI se</span>
        </div>
    </div>
    <div class="tips-card">
        <div style="color:#e2e8f0; font-weight:700; font-size:1rem; margin-bottom:0.8rem;">📌 Quick Study Tips</div>
        <div class="tip-item">⏰ &nbsp; Roz <b style="color:#a78bfa">2-3 ghante</b> focused study karo</div>
        <div class="tip-item">🔁 &nbsp; Galat answers <b style="color:#a78bfa">dobara padho</b></div>
        <div class="tip-item">📝 &nbsp; Important topics ke <b style="color:#a78bfa">short notes</b> banao</div>
        <div class="tip-item">🎯 &nbsp; Ek din mein <b style="color:#a78bfa">ek topic</b> master karo</div>
        <div class="tip-item">😴 &nbsp; <b style="color:#a78bfa">7-8 ghante</b> neend lena zaruri hai</div>
    </div>
    """, unsafe_allow_html=True)

# GENERATE
if uploaded_file and generate_btn:
    st.session_state.submitted = False
    st.session_state.user_answers = {}
    st.session_state.quote = random.choice(QUOTES)

    with st.spinner("🤖 AI padh raha hai..."):
        text = read_uploaded_file(uploaded_file)

    if text:
        topics = extract_topics(text)

        st.markdown('<div class="divider"></div>', unsafe_allow_html=True)
        st.markdown('<div class="section-title">📋 Topics Mile</div>', unsafe_allow_html=True)
        badges = "".join([f'<span class="badge">{t}</span>' for t in topics[:15]])
        st.markdown(f'<div class="glass-card">{badges}</div>', unsafe_allow_html=True)

        all_questions = []
        for topic in topics[:3]:
            with st.spinner(f"Generating {topic} quiz..."):
                quiz_text = generate_quiz(topic)
            qs = parse_quiz(quiz_text)
            for q in qs:
                q["topic"] = topic
            all_questions.extend(qs)

        st.session_state.questions_data = all_questions

# QUIZ
if st.session_state.questions_data:
    st.markdown('<div class="divider"></div>', unsafe_allow_html=True)
    st.markdown('<div class="section-title">❓ Quiz Attempt Karo</div>', unsafe_allow_html=True)

    for i, q in enumerate(st.session_state.questions_data):
        question_text = q.get("question", "")
        options = q.get("options", [])

        st.markdown(f'''
        <div class="q-card">
            <div class="q-number">Question {i+1} • {q.get("topic","")}</div>
            <div class="q-text">{question_text}</div>
        </div>
        ''', unsafe_allow_html=True)

        if options:
            option_labels = [opt[2:].strip() for opt in options]
            option_keys = [opt[0] for opt in options]

            selected = st.radio(
                f"Q{i+1}",
                options=option_keys,
                format_func=lambda x, labels=dict(zip(option_keys, option_labels)): f"{x}. {labels[x]}",
                key=f"q_{i}",
                index=None,
                label_visibility="collapsed"
            )
            st.session_state.user_answers[i] = selected

    st.markdown("<br>", unsafe_allow_html=True)
    submit_btn = st.button("🎯 Submit Quiz & Check Score")

    if submit_btn or st.session_state.submitted:
        st.session_state.submitted = True
        correct_count = 0
        total = len(st.session_state.questions_data)

        st.markdown('<div class="divider"></div>', unsafe_allow_html=True)
        st.markdown('<div class="section-title">📊 Results</div>', unsafe_allow_html=True)

        for i, q in enumerate(st.session_state.questions_data):
            correct = q.get("answer", "").strip().upper()
            user_ans = st.session_state.user_answers.get(i, "") or ""
            is_correct = user_ans.upper() == correct

            if is_correct:
                correct_count += 1
                status = "✅ Sahi!"
                card_class = "correct"
            else:
                status = f"❌ Galat! Sahi Answer: {correct}"
                card_class = "wrong"

            st.markdown(f'''
            <div class="q-card {card_class}">
                <div class="q-number">Question {i+1}</div>
                <div class="q-text">{q.get("question","")}</div>
                <div style="font-size:0.9rem; margin-top:0.5rem;">
                    Tumhara Answer: <b>{user_ans}</b> &nbsp;|&nbsp; {status}
                </div>
            </div>
            ''', unsafe_allow_html=True)

        percentage = int((correct_count / total) * 100) if total > 0 else 0

        if percentage >= 80:
            emoji = "🏆"
            msg = "Excellent!"
        elif percentage >= 60:
            emoji = "👍"
            msg = "Good Job!"
        elif percentage >= 40:
            emoji = "📚"
            msg = "Aur Padho!"
        else:
            emoji = "💪"
            msg = "Keep Trying!"

        st.markdown(f'''
        <div class="score-box">
            <div style="font-size:3rem;">{emoji}</div>
            <div class="score-num">{correct_count}/{total}</div>
            <div style="color:#a78bfa; font-size:1.5rem; font-weight:700;">{percentage}% — {msg}</div>
        </div>
        ''', unsafe_allow_html=True)

# CHATBOT
st.markdown('<div class="divider"></div>', unsafe_allow_html=True)
st.markdown('<div class="section-title">💬 AI Doubt Solver</div>', unsafe_allow_html=True)
st.markdown('<div style="color:#64748b; font-size:0.9rem; margin-bottom:1rem;">Koi bhi subject ka doubt poochho — AI turant jawab dega!</div>', unsafe_allow_html=True)

# Chat history display
for chat in st.session_state.chat_history:
    if chat["role"] == "user":
        st.markdown(f'''
        <div class="chat-user">
        🧑 &nbsp; {chat["content"]}
        </div>
        ''', unsafe_allow_html=True)
    else:
        st.markdown(f'''
        <div class="chat-ai">
        🤖 &nbsp; {chat["content"]}
        </div>
        ''', unsafe_allow_html=True)

# Chat input
user_question = st.chat_input("Ask Any Question.....")

if user_question:
    st.session_state.chat_history.append({
        "role": "user",
        "content": user_question
    })

    with st.spinner("AI soch raha hai..."):
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are a helpful study assistant for students. Answer questions clearly and concisely. Reply in the same language the user uses."
                }
            ] + st.session_state.chat_history
        )
        ai_reply = response.choices[0].message.content

    st.session_state.chat_history.append({
        "role": "assistant",
        "content": ai_reply
    })

    st.rerun()