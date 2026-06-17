import fitz  # pymupdf
from pptx import Presentation

def read_pdf(file_path):
    text = ""
    doc = fitz.open(file_path)
    print(f"PDF - Total pages: {len(doc)}")
    for i, page in enumerate(doc):
        page_text = page.get_text()
        if page_text:
            text += page_text
    return text

def read_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()
    print(f"TXT - Text mila!")
    return text

def read_ppt(file_path):
    text = ""
    prs = Presentation(file_path)
    print(f"PPT - Total slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + "\n"
    return text

def extract_text(file_path):
    if file_path.endswith(".pdf"):
        return read_pdf(file_path)
    elif file_path.endswith(".txt"):
        return read_txt(file_path)
    elif file_path.endswith(".pptx"):
        return read_ppt(file_path)
    else:
        print("❌ Format support nahi hai!")
        return ""

# Test
file = "notes.txt"  # yahan apna file naam likho
text = extract_text(file)

if text:
    print("\n✅ Text mila:")
    print(text)
else:
    print("❌ Text nahi mila!")